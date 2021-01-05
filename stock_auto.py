
import pandas as pd


#################################
## 함수 정의
#################################
def get_stock_code():
    # 종목코드 다운로드
    stock_code = pd.read_html('http://kind.krx.co.kr/corpgeneral/corpList.do?method=download', header=0)[0]

    # 필요없는 column들은 제외
    stock_code = stock_code[['회사명', '종목코드']]

    # 한글 컬럼명을 영어로 변경
    stock_code = stock_code.rename(columns={'회사명': 'company', '종목코드': 'code'})

    # 종목코드가 6자리이기 때문에 6자리를 맞춰주기 위해 설정해줌
    stock_code.code = stock_code.code.map('{:06d}'.format)

    return stock_code


def get_stock(code):
    df = pd.DataFrame()
    for page in range(1, 21):
        # 일별 시세 url
        url = 'http://finance.naver.com/item/sise_day.nhn?code={code}'.format(code=code)
        url = '{url}&page={page}'.format(url=url, page=page)
        print(url)
        current_df = pd.read_html(url, header=0)[0]
        df = df.append(current_df, ignore_index=True)

    return df


def clean_data(df):
    # df.dropna()를 이용해 결측값 있는 행 제거
    df = df.dropna()

    # 한글로 된 컬럼명을 영어로 바꿔줌
    df = df.rename(
        columns={'날짜': 'date', '종가': 'close', '전일비': 'diff', '시가': 'open', '고가': 'high', '저가': 'low', '거래량': 'volume'})
    # 데이터의 타입을 int형으로 바꿔줌
    df[['close', 'diff', 'open', 'high', 'low', 'volume']] = df[
        ['close', 'diff', 'open', 'high', 'low', 'volume']].astype(int)

    # 컬럼명 'date'의 타입을 date로 바꿔줌
    df['date'] = pd.to_datetime(df['date'])

    # 일자(date)를 기준으로 오름차순 정렬
    df = df.sort_values(by=['date'], ascending=True)

    return df


#################################
## 함수 호출
#################################
# 종목 코드 가져오기
company = '삼성전자'
stock_code = get_stock_code()

# 일별 시세 가져오기
code = stock_code[stock_code.company == company].code.values[0].strip()  ## strip() : 공백제거
df = get_stock(code)

# 일별 시세 클린징
df = clean_data(df)

import matplotlib.pyplot as plt
from pandas.plotting import table
import os

# %matplotlib inline 은 jupyter notebook 사용자 용 - jupyter notebook 내에 그래프가 그려지게 한다.

#################################
## 차트 그리기
#################################
plt.figure(figsize=(10, 4))
plt.plot(df['date'], df['close'])
plt.xlabel('date')
plt.ylabel('close')

#################################
## 차트 저장 및 출력하기
#################################
chart_fname = os.path.join("res/stock_report", '{company}_chart.png'.format(company=company))
plt.savefig(chart_fname)
plt.show()

#################################
## 일별 시세 그리기
#################################
plt.figure(figsize=(15, 4))
ax = plt.subplot(111, frame_on=False)  # no visible frame
ax.xaxis.set_visible(False)  # hide the x axis
ax.yaxis.set_visible(False)  # hide the y axis
df = df.sort_values(by=['date'], ascending=False)
table(ax, df.head(10), loc='center', cellLoc='center', rowLoc='center')  # where df is your data frame

#################################
## 일별 시세 저장하기
#################################
table_fname = os.path.join("res/stock_report", '{company}_table.png'.format(company=company))
plt.savefig(table_fname)

import datetime
from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해
import os

#################################
## 파워포인트 객체 선언
#################################
today = datetime.datetime.today().strftime('%Y%m%d')
prs = Presentation()  # 파워포인트 객체 선언

#################################
## 제목 장표 추가
#################################
title_slide_layout = prs.slide_layouts[0]  # 0 : 제목슬라이드에 해당
slide = prs.slides.add_slide(title_slide_layout)  # 제목 슬라이드를 파워포인트 객체에 추가

# 제목 - 제목에 값넣기
title = slide.shapes.title  # 제목
title.text = "주식 보고서"  # 제목에 값 넣기
# 부제목
subtitle = slide.placeholders[1]  # 제목상자는 placeholders[0], 부제목상자는 [1]
subtitle.text = "보고서 작성일 : {date}".format(date=today)

#################################
## 차트 및 테이블 장표 추가
#################################
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

shapes = slide.shapes
shapes.title.text = '{company}, {close}원에 거래 마감'.format(company=company, close=df.iloc[0]['close'])
print(shapes.title.text)

# 차트 추가
left = Inches(0.5)
height = Inches(2.5)
width = Inches(9)
top = Inches(2)
# width, hegith가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(chart_fname, left, top, width=width, height=height)

# 테이블 추가
left = Inches(-1)
height = Inches(3)
width = Inches(12)
top = Inches(4)

pic = slide.shapes.add_picture(table_fname, left, top, width=width, height=height)
cursor_sp = slide.shapes[0]._element
cursor_sp.addprevious(pic._element)  # 해당 요소를 뒤로 보내기 합니다.

#################################
## 보고서 저장
#################################
ppt_fname = os.path.join("res/stock_report", 'stock_report.pptx')
prs.save(ppt_fname)
