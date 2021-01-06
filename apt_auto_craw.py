login = {
    "id": "<네이버 로그인 아이디>",
    "pw": "<네이버 로그인 비밀번호>"
}

# 자동화 테스트를 위해 셀리니움을 불러옵니다.
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
import time

# 접속할 url
url = "https://aptgin.com/"

options = Options()
options.add_argument("--start-maximized")  ## 최대사이즈로 열기
# 크롬 웹 드라이버의 경로를 설정합니다. ex) "C:\chrome_driver\80\chromedriver.exe"
driver = webdriver.Chrome("<드라이버 경로 설정>", chrome_options=options)

# 접속 시도
driver.get(url)
time.sleep(5)

from selenium.webdriver.common.action_chains import ActionChains
from selenium.webdriver.common.keys import Keys
import pyperclip


def clipboard_input(user_xpath, user_input):
    temp_user_input = pyperclip.paste()  # 사용자 클립보드를 따로 저장

    pyperclip.copy(user_input)
    driver.find_element_by_xpath(user_xpath).click()
    ActionChains(driver).key_down(Keys.CONTROL).send_keys('v').key_up(Keys.CONTROL).perform()

    pyperclip.copy(temp_user_input)  # 사용자 클립보드에 저장 된 내용을 다시 가져 옴
    time.sleep(1)


def aptgin_login():
    # 로그인 버튼 누르기
    driver.find_element_by_xpath('//*[@id="gin-header"]/div/div[3]/ul/li[1]/a').click()
    time.sleep(2)
    driver.find_element_by_xpath('//*[@id="naver-login-btn"]').click()
    time.sleep(2)

    ## id/pw 입력
    clipboard_input('//*[@id="id"]', login.get("id"))
    clipboard_input('//*[@id="pw"]', login.get("pw"))
    ## 로그인 버튼 클릭
    driver.find_element_by_xpath('//*[@id="log.login"]').click()
    time.sleep(2)
    # 브라우저 등록 클릭
    driver.find_element_by_xpath('//*[@id="new.save"]').click()


aptgin_login()

import pandas as pd

code_file = "법정동코드 전체자료/법정동코드 전체자료.txt"
code = pd.read_csv(code_file, sep='\t')
code.columns = ['code', 'name', 'is_exist']
code = code[code['is_exist'] == '존재']
## 시, 구 까지만 필요한 상황이므로, 시/구를 분리한다.
addr = code['name'].str.split(" ", n=3, expand=True)

## 컬럼 이름을 변경한다.
addr.columns = ['depth0', 'depth1', 'depth2', 'depth3']

# 시, 구가 중복되어 있으므로, 중복되지 않기 위해 2개의 depth를 연결한후
addr['sigu'] = addr['depth0'] + " " + addr['depth1']
# unique한 값을 알아낸다. 0번째는 nan이 있다. 그래서 1번째 index부터 요청하도록 하겠다.
sigu_list = list(addr.sigu.unique())

len(sigu_list)  # 총 253개의 시+구 가 있나보다.

si = "서울특별시"
gu = "성북구"

# 수요/ 입주 버튼 클릭
driver.find_element_by_xpath('//*[@id="gin-menu"]/ul/li[6]').click()
# 로딩이 좀 오래 걸리므로
time.sleep(10)

driver.find_element_by_xpath('//*[@id="loc1"]').send_keys(si)
time.sleep(2)
driver.find_element_by_xpath('//*[@id="loc2"]').send_keys(gu)
time.sleep(5)
driver.find_element_by_xpath('//*[@id="topSearchForm"]/div/div[6]/button').click()
time.sleep(2)

# 그래프 저장
driver.execute_script("arguments[0].scrollIntoView();", \
                      driver.find_element_by_css_selector(
                          'div.gin-filter'))  # driver.find_element_by_xpath('/html/body/div[5]/div/div[2]') 와 동일코드
time.sleep(2)
div = driver.find_element_by_css_selector(
    'div#chartDivT1')  # div = driver.find_element_by_xpath('//*[@id="chartDivT1"]') 와 동일코드
div.screenshot("graph\%s_%s_수공급.png" % (si, gu))

# 아파트 버튼 클릭
driver.find_element_by_xpath('//*[@id="gin-menu"]/ul/li[2]/a').click()
# 로딩이 좀 오래 걸리므로
time.sleep(10)

el = driver.find_element_by_xpath('//*[@id="serachForm"]/div/div[1]')
######################
## 시 설정
#####################
el.find_element_by_id("navLoc1").click()
ops = el.find_elements_by_tag_name('li')
for op in ops:
    if op.text == si:
        print(op.text)
        op.click()
        break
######################
## 구 설정
#####################
el.find_element_by_id("navLoc2").click()
ops = el.find_elements_by_tag_name('li')
for op in ops:
    if op.text == gu:
        print(op.text)
        op.click()
        break

time.sleep(2)

######################
## 평당 버튼 클릭
#####################
driver.find_element_by_xpath('//*[@id="tab-main"]/div[1]/div[2]/div/label[2]').click()

####################
## 표 저장
####################
driver.execute_script("arguments[0].scrollIntoView();", \
                      driver.find_element_by_css_selector('div.comparer'))
time.sleep(5)
div = driver.find_element_by_css_selector(
    'div#sub02-Apt-Grid')  # div = driver.find_element_by_id('sub02-Apt-Grid') 와 동일코드
div.screenshot("graph\%s_%s_대장아파트.png" % (si, gu))
time.sleep(2)

# 대장아파트로 진입 - 매매전세 가격
driver.find_element_by_xpath('//*[@id="sub02-Apt-Grid"]/div[2]/table/tbody/tr[1]/td[1]/a').click()
time.sleep(5)

driver.execute_script("arguments[0].scrollIntoView();", \
                      driver.find_element_by_css_selector('div.title-area.blue'))
time.sleep(5)
div = driver.find_element_by_css_selector('div.highcharts-container')
div.screenshot("graph\%s_%s_대장아파트매매전세.png" % (si, gu))
time.sleep(2)

######################################
## 수공급 그래프 반복 다운로드
######################################
for sigu in sigu_list[6:10]:
    print(sigu)
    si = sigu.split(" ")[0]
    gu = sigu.split(" ")[1]
    driver.find_element_by_xpath('//*[@id="loc1"]').send_keys(si)
    time.sleep(2)
    driver.find_element_by_xpath('//*[@id="loc2"]').send_keys(gu)
    time.sleep(5)
    driver.find_element_by_xpath('//*[@id="topSearchForm"]/div/div[6]/button').click()
    time.sleep(2)

    ## executeScript("arguments[0].scrollIntoView();",Element );
    ## "arguments[0]" means first index of page starting at 0.
    ## Where an " Element " is the locator on the web page.
    driver.execute_script("arguments[0].scrollIntoView();", \
                          driver.find_element_by_css_selector('div.gin-filter'))
    time.sleep(2)
    div = driver.find_element_by_css_selector('div#chartDivT1')
    div.screenshot("graph\%s_%s_수공급.png" % (si, gu))

for sigu in sigu_list[6:10]:
    # 아파트 버튼 클릭
    driver.find_element_by_xpath('//*[@id="gin-menu"]/ul/li[2]/a').click()
    # 로딩이 좀 오래 걸리므로
    time.sleep(10)

    print(sigu)
    si = sigu.split(" ")[0]
    gu = sigu.split(" ")[1]
    el = driver.find_element_by_xpath('//*[@id="serachForm"]/div/div[1]')

    el.find_element_by_id("navLoc1").click()
    time.sleep(2)
    ops = el.find_elements_by_tag_name('li')
    for op in ops:
        if op.text == si:
            print(op.text)
            op.click()
            break

    time.sleep(2)
    el.find_element_by_id("navLoc2").click()
    time.sleep(2)
    ops = el.find_elements_by_tag_name('li')
    for op in ops:
        if op.text == gu:
            print(op.text)
            op.click()
            break

    time.sleep(2)
    ## 평당 버튼 클릭
    driver.find_element_by_xpath('//*[@id="tab-main"]/div[1]/div[2]/div/label[2]').click()

    # driver.set_window_size(1280, 650)
    driver.execute_script("arguments[0].scrollIntoView();", \
                          driver.find_element_by_css_selector('div.comparer'))
    time.sleep(5)
    div = driver.find_element_by_css_selector('div#sub02-Apt-Grid')
    div.screenshot("graph\%s_%s_대장아파트.png" % (si, gu))
    time.sleep(2)

    # 대장아파트로 진입 - 매매전세 가격
    driver.find_element_by_xpath('//*[@id="sub02-Apt-Grid"]/div[2]/table/tbody/tr[1]/td[1]/a').click()
    time.sleep(5)
    # driver.set_window_size(1280, 1000)
    driver.execute_script("arguments[0].scrollIntoView();", \
                          driver.find_element_by_css_selector('div.title-area.blue'))
    time.sleep(5)
    div = driver.find_element_by_css_selector('div.highcharts-container')
    div.screenshot("graph\%s_%s_대장아파트매매전세.png" % (si, gu))
    time.sleep(2)

## 객체 선언
from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해

prs = Presentation()  # 파워포인트 객체 선언

## 제목슬라이드
title_slide_layout = prs.slide_layouts[0]  # 0: 제목슬라이드
slide = prs.slides.add_slide(title_slide_layout)
# > 제목
title = slide.shapes.title
title.text = "지역별 아파트 분석"
# > 부제목
subtitle = slide.placeholders[1]
subtitle.text = "작성자 : 홍길동"


# 이미지 슬라이드 추가 (함수)
def add_imagelayer(img_pth, title_text):
    img_path = img_pth

    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)

    shapes = slide.shapes
    shapes.title.text = title_text

    left = Inches(0.5)
    height = Inches(5)
    width = Inches(9)
    top = Inches(2)
    # width, hegith가 없을 경우 원본 사이즈로
    pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)


# 반복문으로 이미지 슬라이드 추가하기
for sigu in sigu_list[1:6]:
    si = sigu.split("경기도")[0]
    gu = sigu.split("파주시")[1]
    print(si, gu)
    add_imagelayer('graph\\%s_%s_수공급.png' % (si, gu), '%s %s 수공급' % (si, gu))
    add_imagelayer('graph\\%s_%s_대장아파트.png' % (si, gu), '%s %s 아파트 목록' % (si, gu))
    add_imagelayer('graph\\%s_%s_대장아파트매매전세.png' % (si, gu), '%s %s 대장아파트매매전세' % (si, gu))

# 저장
prs.save('apt_report.pptx')
