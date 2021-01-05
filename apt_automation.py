## 객체 선언
from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해

prs = Presentation()  # 파워포인트 객체 선언

## 제목슬라이드
title_slide_layout = prs.slide_layouts[0]  # 0: 제목슬라이드
slide = prs.slides.add_slide(title_slide_layout)
# > 제목
title = slide.shapes.title
title.text = "지역별 아파트 분석"
# > 부제목
subtitle = slide.placeholders[1]
subtitle.text = "작성자 : 홍길동"


# 이미지 슬라이드 추가 (함수)
def add_imagelayer(img_pth, title_text):
    img_path = img_pth

    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)

    shapes = slide.shapes
    shapes.title.text = title_text

    left = Inches(0.5)
    height = Inches(5)
    width = Inches(9)
    top = Inches(2)
    # width, hegith가 없을 경우 원본 사이즈로
    pic = slide.shapes.add_picture(img_path, left, top, width=width, height=height)


# 반복문으로 이미지 슬라이드 추가하기
for sigu in sigu_list[1:6]:
    si = sigu.split("경기도")[0]
    gu = sigu.split("파주시")[1]
    print(si, gu)
    add_imagelayer('graph\\%s_%s_수공급.png' % (si, gu), '%s %s 수공급' % (si, gu))
    add_imagelayer('graph\\%s_%s_대장아파트.png' % (si, gu), '%s %s 아파트 목록' % (si, gu))
    add_imagelayer('graph\\%s_%s_대장아파트매매전세.png' % (si, gu), '%s %s 대장아파트매매전세' % (si, gu))

# 저장
prs.save('apt_report.pptx')